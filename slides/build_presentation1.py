"""
Presentation 1 (Week 6-7, 10 minutes): define the area/focus of the project,
what we're trying to understand, and how we plan to go about it. No regression
results -- those are Presentation 2's job.

Narrative: puzzle -> our question -> what prior work found (the gap) -> a model
of peer effects -> the model's sharp prediction (effort tracks Katz-Bonacich
centrality) -> what phi means (a dial from popularity to reach) -> predictions
-> the data -> how we estimate it -> which way the arrow runs -> the plan.

Matches research_brief.pdf. Run build_assets.py first to (re)generate images.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "assets"

BLUE = RGBColor(0x2A, 0x78, 0xD6)
ORANGE = RGBColor(0xEB, 0x68, 0x34)
DARK = RGBColor(0x19, 0x18, 0x0F)
GRAY_TEXT = RGBColor(0x55, 0x50, 0x3F)
LIGHT_BG = RGBColor(0xFA, 0xF9, 0xF6)
CARD_BG = RGBColor(0xF1, 0xEF, 0xE8)
BORDER = RGBColor(0xDD, 0xD8, 0xCD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE_TINT = RGBColor(0xE9, 0xF0, 0xFA)

FONT_HEAD = "Georgia"
FONT_BODY = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide(bg_color=LIGHT_BG):
    slide = prs.slides.add_slide(BLANK)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()
    bg.shadow.inherit = False
    spTree = slide.shapes._spTree
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return slide


def add_text(slide, left, top, width, height, text, size=18, bold=False,
             color=DARK, font=FONT_BODY, align=PP_ALIGN.LEFT, italic=False,
             line_spacing=1.15, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = line_spacing
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color
            run.font.name = font
    return box


def add_bullets(slide, left, top, width, height, items, size=16, color=DARK,
                font=FONT_BODY, space_after=10, line_spacing=1.2):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = "  •  " + item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = font
    return box


def add_pagenum(slide, n, color=GRAY_TEXT):
    add_text(slide, Inches(12.5), Inches(7.08), Inches(0.7), Inches(0.35),
             str(n), size=11, color=color, align=PP_ALIGN.RIGHT)


def eyebrow(slide, text, color=ORANGE):
    add_text(slide, Inches(0.7), Inches(0.4), Inches(11), Inches(0.4), text,
             size=13, bold=True, color=color, font=FONT_BODY)


def title(slide, text, top=Inches(0.78), size=27, width=Inches(12.0), color=DARK):
    add_text(slide, Inches(0.7), top, width, Inches(1.0), text, size=size,
             bold=True, color=color, font=FONT_HEAD, line_spacing=1.08)


def rounded_box(slide, left, top, width, height, text="", fill=WHITE,
                text_color=DARK, size=14, bold=False, border=BORDER,
                align=PP_ALIGN.CENTER, font=FONT_BODY):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.12
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(1.25)
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = text_color
        run.font.name = font
    return shape


def arrow(slide, x1, y1, x2, y2, color=GRAY_TEXT, width=Pt(2)):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      int(x1), int(y1), int(x2), int(y2))
    conn.line.color.rgb = color
    conn.line.width = width
    ln = conn.line._get_or_add_ln()
    tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail)
    return conn


def pic(slide, fname, left, top, width=None, height=None):
    kw = {}
    if width is not None:
        kw["width"] = width
    if height is not None:
        kw["height"] = height
    return slide.shapes.add_picture(str(ASSETS / fname), left, top, **kw)


W = Inches(13.333)


def centered_left(width):
    return (W - width) / 2


# ===================================================== 1 -- title
s = add_slide()
pic(s, "decor_network.png", Inches(9.05), Inches(0.35), height=Inches(4.0))
add_text(s, Inches(0.9), Inches(2.05), Inches(11.5), Inches(0.5),
         "ECC3841 NETWORK ECONOMICS  ·  PRESENTATION 1",
         size=13, bold=True, color=ORANGE)
add_text(s, Inches(0.9), Inches(2.65), Inches(11.5), Inches(1.9),
         "Network Position and\nAcademic Performance",
         size=44, bold=True, color=DARK, font=FONT_HEAD, line_spacing=1.05)
add_text(s, Inches(0.9), Inches(4.75), Inches(11.4), Inches(1.0),
         "Does where you sit in a friendship network predict your grades,\n"
         "beyond who your direct friends are?",
         size=18, color=GRAY_TEXT, italic=True, line_spacing=1.3)
add_text(s, Inches(0.9), Inches(6.55), Inches(9), Inches(0.5),
         "Team [names]  ·  [Date]", size=13, color=GRAY_TEXT)

# ===================================================== 2 -- the puzzle
s = add_slide()
eyebrow(s, "THE PUZZLE")
title(s, "Two students. Same-looking friend groups. Same grades?")
pic(s, "student_comparison.png", centered_left(Inches(9.4)), Inches(1.75), width=Inches(9.4))
rounded_box(s, Inches(0.7), Inches(6.2), Inches(11.9), Inches(0.95),
            "Does a student's position in a friendship network affect their grades?",
            fill=BLUE_TINT, border=BLUE, text_color=DARK, size=21, bold=True, font=FONT_HEAD)

# ===================================================== 3 -- prior work / the gap
s = add_slide()
eyebrow(s, "WHAT WE ALREADY KNOW  ·  PRIOR WORK")
title(s, "Smirnov & Thurner (2017): the effect is mostly selection")
add_text(s, Inches(0.7), Inches(1.7), Inches(12.0), Inches(0.5),
         "A real, changing “who-likes-who” network of Russian students, matched to GPA.",
         size=15, italic=True, color=GRAY_TEXT)

y, bw, bh, gap = Inches(2.5), Inches(3.55), Inches(1.05), Inches(0.5)
x0 = Inches(0.85)
labels = ["Students pick friends\nwith similar GPA",
          "Friend groups end up\nGPA-similar",
          "GPA homophily\n(what we observe)"]
xs = []
for i, lab in enumerate(labels):
    x = x0 + i * (bw + gap)
    xs.append(x)
    rounded_box(s, x, y, bw, bh, lab, fill=WHITE, size=14, bold=True)
for i in range(2):
    arrow(s, xs[i] + bw, y + bh / 2, xs[i + 1], y + bh / 2)

add_text(s, Inches(0.85), Inches(3.95), Inches(11.6), Inches(0.6),
         "Their measure is one number per student — the average GPA of their direct friends.",
         size=16, color=GRAY_TEXT)
add_text(s, Inches(0.85), Inches(4.75), Inches(11.6), Inches(0.7),
         "→  It never asks where that student sits in the wider network.",
         size=19, bold=True, color=ORANGE)
add_text(s, Inches(0.85), Inches(5.75), Inches(11.6), Inches(1.0),
         "Two students with the same direct friends score the same — however different "
         "those friends' own social worlds are.",
         size=15, color=GRAY_TEXT, line_spacing=1.3)

# ===================================================== 4 -- the model
s = add_slide()
eyebrow(s, "THE MODEL", color=BLUE)
title(s, "A model of peer effects in study effort")
add_text(s, Inches(0.7), Inches(1.75), Inches(12.0), Inches(0.5),
         "Ballester, Calvó-Armengol & Zenou (2006)  ·  Calvó-Armengol, Patacchini & Zenou (2009; IZA DP 3859)",
         size=12, italic=True, color=GRAY_TEXT)

rounded_box(s, Inches(1.3), Inches(2.35), Inches(10.7), Inches(1.5), fill=CARD_BG)
pic(s, "eq_utility.png", centered_left(Inches(8.1)), Inches(2.62), width=Inches(8.1))

add_text(s, Inches(0.9), Inches(4.25), Inches(11.5), Inches(0.5),
         "Each student i picks study effort xᵢ. Three parts:", size=17, bold=True, color=DARK)
add_bullets(s, Inches(1.1), Inches(4.8), Inches(11.2), Inches(2.4), [
    "a xᵢ  —  what you gain from effort",
    "−½xᵢ²  —  its rising cost: studying alone has diminishing returns",
    "φ Σⱼ gᵢⱼ xᵢ xⱼ  —  the social part: when a friend studies more, your own payoff from "
    "studying goes up  (φ > 0)",
], size=16, color=GRAY_TEXT, space_after=13)

# ===================================================== 5 -- equilibrium
s = add_slide()
eyebrow(s, "THE MODEL", color=BLUE)
title(s, "Everyone best-responds → effort tracks one measure", size=25)

add_text(s, Inches(0.9), Inches(1.85), Inches(2.5), Inches(0.4), "Best response",
         size=14, bold=True, color=GRAY_TEXT)
pic(s, "eq_bestresponse.png", Inches(3.5), Inches(1.68), width=Inches(3.5))
add_text(s, Inches(7.5), Inches(1.85), Inches(5.3), Inches(0.6),
         "study a baseline, plus more when your friends study more",
         size=13, italic=True, color=GRAY_TEXT, line_spacing=1.2)

add_text(s, Inches(0.9), Inches(2.95), Inches(2.5), Inches(0.4), "Solve for all students",
         size=14, bold=True, color=GRAY_TEXT)
pic(s, "eq_equilibrium.png", Inches(3.5), Inches(2.78), width=Inches(5.6))

add_text(s, Inches(0.9), Inches(4.0), Inches(2.5), Inches(0.4), "Valid while",
         size=14, bold=True, color=GRAY_TEXT)
pic(s, "eq_condition.png", Inches(3.5), Inches(3.85), width=Inches(2.15))

rounded_box(s, Inches(0.9), Inches(4.85), Inches(11.5), Inches(2.05),
            fill=BLUE_TINT, border=BLUE)
add_text(s, Inches(1.25), Inches(5.05), Inches(10.9), Inches(1.7),
         "A sharp, falsifiable claim: achievement should track Katz-Bonacich centrality "
         "b(g,φ) — not raw popularity, not any other network measure. φ is the model's own "
         "parameter, and the one we vary in the data. In the original AddHealth study, a "
         "1-SD rise in centrality raised school performance by about 7% of a SD.",
         size=15, color=DARK, line_spacing=1.32)

# ===================================================== 6 -- what phi means
s = add_slide()
eyebrow(s, "THE KEY IDEA", color=ORANGE)
title(s, "φ is a dial — from popularity to reach")
pic(s, "phi_dial.png", centered_left(Inches(11.9)), Inches(1.55), width=Inches(11.9))
rounded_box(s, Inches(0.9), Inches(6.15), Inches(11.5), Inches(1.0),
            "Turn φ up and the same formula stops asking “how many friends?” and starts "
            "asking “where in the network?” — we test GPA against the whole dial.",
            fill=CARD_BG, border=BORDER, size=15, font=FONT_BODY)

# ===================================================== 7 -- predictions
s = add_slide()
eyebrow(s, "PREDICTIONS")
title(s, "What we expect — and what would prove us wrong")

rounded_box(s, Inches(0.7), Inches(1.7), Inches(11.95), Inches(0.46), "FROM THE MODEL",
            fill=BLUE, text_color=WHITE, size=12.5, bold=True, align=PP_ALIGN.LEFT)
add_text(s, Inches(0.75), Inches(2.24), Inches(11.85), Inches(1.0),
         "Katz-Bonacich centrality predicts next-term GPA — even after friend-group GPA and "
         "raw degree are controlled. Falsified if it loses significance once degree is in, "
         "or if it flips sign along the φ dial.",
         size=15, color=DARK, line_spacing=1.28)

add_text(s, Inches(0.75), Inches(3.5), Inches(11), Inches(0.35),
         "IF IT PREDICTS — WHICH DIRECTION?", size=12.5, bold=True, color=GRAY_TEXT)
tw, tgap = Inches(3.83), Inches(0.22)
tx = Inches(0.7)
tri = [("β₁ > 0", "access to help lifts grades"),
       ("β₁ = 0", "help and time-cost cancel"),
       ("β₁ < 0", "wide reach eats study time")]
for i, (h, d) in enumerate(tri):
    x = tx + i * (tw + tgap)
    rounded_box(s, x, Inches(3.9), tw, Inches(1.0), fill=CARD_BG)
    add_text(s, x + Inches(0.2), Inches(4.0), tw - Inches(0.4), Inches(0.4), h,
             size=15, bold=True, color=DARK, font=FONT_HEAD)
    add_text(s, x + Inches(0.2), Inches(4.42), tw - Inches(0.4), Inches(0.5), d,
             size=12.5, color=GRAY_TEXT)

add_text(s, Inches(0.75), Inches(5.2), Inches(11), Inches(0.35),
         "FROM REPLICATING SMIRNOV & THURNER", size=12.5, bold=True, color=GRAY_TEXT)
add_bullets(s, Inches(0.9), Inches(5.6), Inches(11.6), Inches(1.7), [
    "A friend's past GPA won't predict a student's future GPA, once own past GPA is controlled.",
    "New friendships show a smaller GPA gap than ending ones — a sign of selection.",
    "Selection is stronger at university than in high school.",
], size=14, color=GRAY_TEXT, space_after=6)

# ===================================================== 8 -- the data
s = add_slide()
eyebrow(s, "THE DATA")
title(s, "A real, directed friendship network — observed over time")
pic(s, "real_network_school.png", Inches(6.7), Inches(1.7), height=Inches(4.9))

add_text(s, Inches(0.7), Inches(1.85), Inches(5.7), Inches(0.35), "NODES", size=13, bold=True, color=BLUE)
add_text(s, Inches(0.7), Inches(2.2), Inches(5.7), Inches(0.9),
         "Students. 655 in high school; 1,200–1,549 per university year.",
         size=14, color=DARK, line_spacing=1.3)

add_text(s, Inches(0.7), Inches(3.15), Inches(5.7), Inches(0.35), "EDGES", size=13, bold=True, color=ORANGE)
add_text(s, Inches(0.7), Inches(3.5), Inches(5.7), Inches(1.5),
         "Directed “like” ties. i → j means i liked j at least once in a 3-month window. "
         "2–14 snapshots per group, 38 in total. Only 24–27% are mutual.",
         size=14, color=DARK, line_spacing=1.3)

add_text(s, Inches(0.7), Inches(5.15), Inches(5.7), Inches(0.35), "ATTRIBUTES", size=13, bold=True, color=DARK)
add_text(s, Inches(0.7), Inches(5.5), Inches(5.7), Inches(1.2),
         "GPA; in/out-degree, Katz-Bonacich, betweenness, eigenvector; school vs. university.",
         size=14, color=DARK, line_spacing=1.3)

rounded_box(s, Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.6),
            "Same dataset as Smirnov & Thurner — our results compare directly with theirs.",
            fill=CARD_BG, border=BORDER, size=13.5, font=FONT_BODY)

# ===================================================== 9 -- how we estimate it
s = add_slide()
eyebrow(s, "HOW WE'LL GO ABOUT IT")
title(s, "Predict next-term GPA from position — holding the rest fixed")

rounded_box(s, Inches(0.7), Inches(1.8), Inches(11.95), Inches(1.9), fill=CARD_BG)
pic(s, "eq_regression.png", centered_left(Inches(11.0)), Inches(2.05), width=Inches(11.0))

add_text(s, Inches(0.9), Inches(4.05), Inches(11.5), Inches(0.45),
         "Each control has a job:", size=17, bold=True, color=DARK)
add_bullets(s, Inches(1.1), Inches(4.6), Inches(11.4), Inches(1.6), [
    "β₃  average friend GPA  →  β₁ is net of friend-group composition (Smirnov & Thurner's own variable)",
    "β₄, β₅  raw in / out-degree  →  β₁ is net of popularity",
], size=15, color=GRAY_TEXT, space_after=10)

add_text(s, Inches(0.9), Inches(6.05), Inches(11.5), Inches(0.8),
         "The real question: does centrality carry anything beyond simply counting friends?",
         size=17, bold=True, color=ORANGE)

# ===================================================== 10 -- which way does the arrow run
s = add_slide()
eyebrow(s, "IDENTIFICATION")
title(s, "Which way does the arrow run?")
add_text(s, Inches(0.7), Inches(1.62), Inches(12.0), Inches(0.75),
         "Does position lift grades — or do good grades just attract friends? "
         "One snapshot can't tell. Our panel data help three ways:",
         size=15, italic=True, color=GRAY_TEXT, line_spacing=1.25)
pic(s, "timeline_snapshots.png", centered_left(Inches(9.0)), Inches(2.35), width=Inches(9.0))

add_bullets(s, Inches(1.0), Inches(4.75), Inches(11.4), Inches(1.7), [
    "Time order — position is measured before the GPA it predicts.",
    "Past GPA controlled — we compare students at the same current grade.",
    "Reverse arrow measured directly — does a GPA rise bring new friends? (Smirnov & Thurner's tie-formation check)",
], size=14.5, color=DARK, space_after=8)

rounded_box(s, Inches(0.7), Inches(6.5), Inches(11.95), Inches(0.62),
            "So our claim is that position predicts later grades — not that it causes them.",
            fill=BLUE_TINT, border=BLUE, size=14.5, bold=True, font=FONT_BODY)

# ===================================================== 11 -- the plan
s = add_slide()
eyebrow(s, "OUR PLAN")
title(s, "Three layers of analysis")

y, bw, bh, gap = Inches(2.1), Inches(3.63), Inches(1.35), Inches(0.5)
x0 = Inches(0.7)
layers = [
    ("LAYER 1", "Replicate", "Rebuild Smirnov & Thurner's\nselection-vs-influence checks\non the same data"),
    ("LAYER 2", "Test the model", "Compute Katz-Bonacich\ncentrality; test the prediction\nacross the φ dial"),
    ("LAYER 3", "Compare", "High school vs. university —\ndoes the effect hold across\nsettings?"),
]
xs = []
for i, (tag, head, _) in enumerate(layers):
    x = x0 + i * (bw + gap)
    xs.append(x)
    rounded_box(s, x, y, bw, bh, fill=WHITE)
    add_text(s, x + Inches(0.25), y + Inches(0.16), bw - Inches(0.5), Inches(0.3), tag,
             size=12.5, bold=True, color=ORANGE)
    add_text(s, x + Inches(0.25), y + Inches(0.5), bw - Inches(0.5), Inches(0.5), head,
             size=19, bold=True, color=DARK, font=FONT_HEAD)
for i in range(2):
    arrow(s, xs[i] + bw, y + bh / 2, xs[i + 1], y + bh / 2)
for i, (_, _, desc) in enumerate(layers):
    add_text(s, xs[i] + Inches(0.05), y + bh + Inches(0.3), bw, Inches(1.6), desc,
             size=13.5, color=GRAY_TEXT, line_spacing=1.3)

add_text(s, Inches(0.7), Inches(5.7), Inches(11.9), Inches(0.8),
         "One headline test, at φ = 0.85 × (1/λmax), is fixed in advance — so later "
         "robustness checks can't be accused of cherry-picking.",
         size=14.5, italic=True, color=GRAY_TEXT, line_spacing=1.25)

rounded_box(s, Inches(0.7), Inches(6.55), Inches(11.95), Inches(0.6),
            "Also grounded in: Calvó-Armengol, Patacchini & Zenou (2009)  ·  "
            "Giulietti, Vlassopoulos & Zenou (2020) on peer depression",
            fill=CARD_BG, border=BORDER, size=12.5, font=FONT_BODY)

# ===================================================== 12 -- questions
s = add_slide()
add_text(s, Inches(0.9), Inches(2.7), Inches(11), Inches(1.2), "Questions?",
         size=46, bold=True, color=DARK, font=FONT_HEAD)
add_bullets(s, Inches(0.95), Inches(4.15), Inches(11.5), Inches(2.0), [
    "Data — Smirnov & Thurner (2017), Harvard Dataverse  doi:10.7910/DVN/SZA9YW",
    "Model — Ballester, Calvó-Armengol & Zenou (2006); Calvó-Armengol, Patacchini & Zenou (2009)",
    "Full detail — research_brief.pdf",
], size=14, color=GRAY_TEXT, space_after=8)

for i, sl in enumerate(prs.slides, start=1):
    add_pagenum(sl, i)

out_path = ROOT / "Presentation1.pptx"
prs.save(out_path)
print(f"Saved -> {out_path}  ({len(prs.slides._sldIdLst)} slides)")
